from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "ChatPilot: Conversational Autonomous Driving Agent"
slide.placeholders[1].text = "Bridging human conversation with autonomous vehicle navigation\nMilestone Presentation"

def add_slide(title, bullets):
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
    return slide

add_slide("Project Overview", [
    "Enables natural language communication with an autonomous rover",
    "Integrates Speech, NLP, and Navigation",
    "Operates in a simulated & real-time driving environment",
    "Objective: Human–AI interaction for autonomous mobility"
])

add_slide("System Architecture", [
    "Mobile App: Captures speech, processes NLP, sends commands",
    "MQTT Broker: Manages communication between app & rover",
    "Rover (Jetson Xavier): Handles navigation & obstacle detection",
    "Sensors: RPLidar for mapping and safety",
    "Visual: Architecture diagram (center)"
])

add_slide("Technology Stack", [
    "Programming: Python",
    "Libraries: PyAudio, pyttsx3, paho-MQTT, DroneKit, MAVProxy",
    "Hardware: Jetson Xavier, RPLidar, Rover platform",
    "Tools: VS Code, GitHub",
    "Visual: Tech logos or stack icons"
])

add_slide("Mathematical Foundation", [
    "Haversine Formula: Calculates distance between coordinates",
    "A-Star Algorithm: Finds shortest safe route",
    "Enables accurate, efficient path planning",
    "Visual: Small formula + path diagram"
])

add_slide("Live Demo", [
    "User gives voice command via mobile app",
    "Command → Intent → Coordinates → Rover navigation",
    "Real-time Lidar-based obstacle avoidance",
    "Feedback displayed on mobile interface",
    "Visuals: Screenshots of app and demo video"
])

add_slide("Results & Learnings", [
    "Successful integration of Speech + NLP + Navigation",
    "Handled latency, Lidar noise, real-time control",
    "Limitations: Response delay, limited multi-turn dialogue",
    "Future Work: Integrate video-language foundation models for better perception",
    "Visual: Rover image or demo summary chart"
])

prs.save("ChatPilot_Milestone_Presentation.pptx")
print("✅ Presentation saved as ChatPilot_Milestone_Presentation.pptx")
